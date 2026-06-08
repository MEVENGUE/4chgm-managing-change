"""Export engine — PDF reports and PPT summaries."""

from __future__ import annotations

import io
from datetime import datetime, timezone

from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

from app.services import storage


def generate_pdf_report(title: str, sections: list[dict]) -> bytes:
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=A4)
    styles = getSampleStyleSheet()
    story = [Paragraph(title, styles["Title"]), Spacer(1, 12)]
    story.append(Paragraph(f"Generated {datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')}", styles["Normal"]))
    story.append(Spacer(1, 18))
    for sec in sections:
        story.append(Paragraph(sec.get("heading", "Section"), styles["Heading2"]))
        story.append(Paragraph(sec.get("body", "").replace("\n", "<br/>"), styles["Normal"]))
        story.append(Spacer(1, 12))
    doc.build(story)
    return buffer.getvalue()


def generate_ppt_summary(title: str, bullets: list[str]) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    layout = prs.slide_layouts[1]
    slide2 = prs.slides.add_slide(layout)
    slide2.shapes.title.text = "Executive Summary"
    body = slide2.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets[:8]):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.font.size = Pt(14)
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def export_and_store(
    file_bytes: bytes,
    file_name: str,
    bucket_type: str,
    organization_id: str,
    workspace_id: str,
    user_id: str,
    content_type: str,
) -> tuple[str, str]:
    return storage.upload_file(
        file_bytes, file_name, bucket_type, organization_id, workspace_id, user_id, content_type
    )
